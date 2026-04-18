# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUTPUTS = Path('outputs')

C_BG       = RGBColor(0xF0, 0xF6, 0xFF)
C_NAVY     = RGBColor(0x0A, 0x2A, 0x5C)
C_BLUE     = RGBColor(0x1A, 0x6F, 0xC8)
C_SKY      = RGBColor(0x41, 0xB3, 0xF0)
C_LIGHT    = RGBColor(0xA8, 0xD8, 0xF7)
C_PALE     = RGBColor(0xE3, 0xF2, 0xFD)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY     = RGBColor(0x6A, 0x7E, 0x9A)
C_DARK     = RGBColor(0x0D, 0x1F, 0x3C)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── 헬퍼 ──────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background; bg.fill.solid(); bg.fill.fore_color.rgb = C_BG
    return s

def box(sl, x, y, w, h, fill=C_WHITE, line=None, lw=Pt(1.2), rnd=False):
    s = sl.shapes.add_shape(5 if rnd else 1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else:    s.line.fill.background()
    if rnd and hasattr(s, 'adjustments'): s.adjustments[0] = 0.07
    return s

def txt(sl, text, x, y, w, h,
        size=Pt(10), color=C_DARK, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = size; r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    r.font.name = 'Malgun Gothic'
    return tb

def img(sl, fname, x, y, w, h):
    return sl.shapes.add_picture(str(OUTPUTS / fname), x, y, w, h)

def title_bar(sl, title, sub):
    box(sl, Inches(0), Inches(0), W, Inches(0.75), fill=C_NAVY)
    box(sl, Inches(0), Inches(0), Inches(0.1), Inches(0.75), fill=C_SKY)
    txt(sl, title, Inches(0.25), Inches(0.1), Inches(8), Inches(0.55),
        size=Pt(22), color=C_WHITE, bold=True)
    txt(sl, sub, Inches(8.3), Inches(0.22), Inches(4.8), Inches(0.35),
        size=Pt(9.5), color=C_LIGHT, align=PP_ALIGN.RIGHT)

def insight_card(sl, number, title, body, x, y, w=Inches(3.6), h=Inches(1.5)):
    """분석 인사이트 카드 — 번호 + 제목 + 본문"""
    box(sl, x, y, w, h, fill=C_WHITE, line=C_LIGHT, lw=Pt(1.5), rnd=True)
    # 번호 원형
    box(sl, x + Inches(0.15), y + Inches(0.15),
        Inches(0.42), Inches(0.42), fill=C_BLUE, rnd=True)
    txt(sl, number,
        x + Inches(0.15), y + Inches(0.15),
        Inches(0.42), Inches(0.42),
        size=Pt(13), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, title,
        x + Inches(0.68), y + Inches(0.15),
        w - Inches(0.82), Inches(0.38),
        size=Pt(10.5), color=C_BLUE, bold=True)
    txt(sl, body,
        x + Inches(0.15), y + Inches(0.62),
        w - Inches(0.3), Inches(0.8),
        size=Pt(9), color=C_GRAY)

def arrow_card(sl, from_txt, to_txt, x, y, w=Inches(3.9)):
    """발견 → 결정 화살표 카드"""
    box(sl, x, y, w, Inches(1.2), fill=C_WHITE, line=C_LIGHT, lw=Pt(1.2), rnd=True)
    box(sl, x, y, w, Inches(0.06), fill=C_SKY)
    txt(sl, from_txt, x + Inches(0.15), y + Inches(0.1),
        w - Inches(0.3), Inches(0.42),
        size=Pt(9), color=C_GRAY)
    txt(sl, '▼',
        x + Inches(0), y + Inches(0.5),
        w, Inches(0.28),
        size=Pt(10), color=C_SKY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, to_txt, x + Inches(0.15), y + Inches(0.76),
        w - Inches(0.3), Inches(0.38),
        size=Pt(9.5), color=C_BLUE, bold=True)


# ══════════════════════════════════════════════════════════════
# S1. 데이터 특성 분석 — 시계열·정상성·계절성
# ══════════════════════════════════════════════════════════════
s1 = new_slide()
title_bar(s1, '데이터 특성 분석',
          '시계열 구조 · 정상성 검정 · 계절성 분해')

# ── 좌: STL 분해 (메인) ──────────────────────────────────
box(s1, Inches(0.25), Inches(0.88), Inches(5.6), Inches(4.55),
    fill=C_WHITE, line=C_LIGHT, lw=Pt(1.2), rnd=True)
txt(s1, 'STL 분해  —  추세 · 계절성 · 잔차',
    Inches(0.4), Inches(0.96), Inches(5.2), Inches(0.34),
    size=Pt(10), color=C_BLUE, bold=True)
img(s1, '05_stl_decomposition.png',
    Inches(0.3), Inches(1.32), Inches(5.5), Inches(3.8))

# ── 우상: ACF/PACF ────────────────────────────────────────
box(s1, Inches(6.05), Inches(0.88), Inches(7.0), Inches(2.25),
    fill=C_WHITE, line=C_LIGHT, lw=Pt(1.2), rnd=True)
txt(s1, 'ACF / PACF  —  자기상관 구조',
    Inches(6.2), Inches(0.96), Inches(6.6), Inches(0.32),
    size=Pt(10), color=C_BLUE, bold=True)
img(s1, '05_acf_pacf.png',
    Inches(6.1), Inches(1.3), Inches(6.85), Inches(1.7))

# ── 우하: 분포 ────────────────────────────────────────────
box(s1, Inches(6.05), Inches(3.22), Inches(7.0), Inches(2.2),
    fill=C_WHITE, line=C_LIGHT, lw=Pt(1.2), rnd=True)
txt(s1, '가격 변화율 분포  —  정규성 검정',
    Inches(6.2), Inches(3.3), Inches(6.6), Inches(0.32),
    size=Pt(10), color=C_BLUE, bold=True)
img(s1, '04_distribution.png',
    Inches(6.1), Inches(3.62), Inches(6.85), Inches(1.68))

# ── 하단 인사이트 3개 카드 ───────────────────────────────
insight_card(s1,
    '01', '비정상(Non-stationary) 시계열',
    '추세·분산이 시간에 따라 변함\n2020년 이후 레짐 전환 구간 존재',
    Inches(0.25), Inches(5.55), w=Inches(3.95))

insight_card(s1,
    '02', 'AR(1) 구조 + 강한 자기상관',
    'ACF 느린 감소 → 단기 자기상관 지배\nPACF 1차 절단 → 단순 AR 패턴',
    Inches(4.4), Inches(5.55), w=Inches(3.95))

insight_card(s1,
    '03', '뚜렷한 계절성·레짐 전환',
    'STL 계절 성분 확인\n급등락 이벤트 → 선형 모델 한계',
    Inches(8.55), Inches(5.55), w=Inches(4.5))

# 하단 결론 바
box(s1, Inches(0), Inches(7.08), W, Inches(0.04), fill=C_SKY)
txt(s1,
    '→  비정상·레짐 전환 특성상 단순 회귀 예측보다  방향(상승/하락) 분류가 더 안정적 접근',
    Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3),
    size=Pt(9.5), color=C_BLUE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# S2. 피처 분석 — 상관관계 · 선행지표 · 중요도
# ══════════════════════════════════════════════════════════════
s2 = new_slide()
title_bar(s2, '피처 분석  —  상관관계 · Lag · 중요도',
          '155종 피처 → 분석 기반 79개 선택')

# ── 좌: Lag 상관관계 ─────────────────────────────────────
box(s2, Inches(0.25), Inches(0.88), Inches(5.85), Inches(4.35),
    fill=C_WHITE, line=C_LIGHT, lw=Pt(1.2), rnd=True)
txt(s2, 'Lag 상관관계  —  피처 선행성 분석',
    Inches(0.4), Inches(0.96), Inches(5.4), Inches(0.32),
    size=Pt(10), color=C_BLUE, bold=True)
img(s2, '06_lag_correlation.png',
    Inches(0.3), Inches(1.3), Inches(5.75), Inches(3.68))

# ── 우: 앙상블 피처 중요도 ───────────────────────────────
box(s2, Inches(6.3), Inches(0.88), Inches(6.8), Inches(4.35),
    fill=C_WHITE, line=C_LIGHT, lw=Pt(1.2), rnd=True)
txt(s2, '피처 중요도  —  MI + XGBoost + LASSO 앙상블',
    Inches(6.45), Inches(0.96), Inches(6.4), Inches(0.32),
    size=Pt(10), color=C_BLUE, bold=True)
img(s2, '07_ensemble_importance.png',
    Inches(6.35), Inches(1.3), Inches(6.7), Inches(3.68))

# ── 하단 인사이트 + 모델링 방향 결정 ─────────────────────
# 인사이트 → 결정 3세트
findings = [
    ('고철 수입가 10~20일 선행\n단가와 Spearman r > 0.8',
     '→  Lag 피처 + 방향성 변화율 피처 설계'),
    ('거시경제·에너지 다중공선성 높음\nVIF 10 이상 다수',
     '→  앙상블 선택으로 중복 피처 제거'),
    ('뉴스심리지수(NSI) 단기 유효\n모멘텀 크로스 신호 상위 랭크',
     '→  MA크로스·상승비율 등 방향성 피처 18개 추가'),
]
fx = Inches(0.25)
for from_t, to_t in findings:
    arrow_card(s2, from_t, to_t, fx, Inches(5.35), w=Inches(4.2))
    fx += Inches(4.38)

# 하단 결론 바
box(s2, Inches(0), Inches(7.08), W, Inches(0.04), fill=C_SKY)
txt(s2,
    '→  선행지표 발굴 + 다중공선성 제거 + 방향성 피처 추가  →  155종에서 79 + 18개 = 97개 최종 피처셋 구성',
    Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3),
    size=Pt(9.5), color=C_BLUE, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# S3. 분석 기반 모델링 결정 → 결과
# ══════════════════════════════════════════════════════════════
s3 = new_slide()
title_bar(s3, '분석 기반 모델링 결정  →  결과',
          '분석 인사이트가 각 결정의 근거')

# ── 상단: 분석 → 결정 → 결과 플로우 ─────────────────────
flow_items = [
    (C_SKY,  '분석 발견',
     '비정상 시계열\n레짐 전환 존재\nAR(1) 구조'),
    (C_BLUE, '모델 방향 결정',
     '회귀 → 방향 분류\n단기(10일) 집중\nWalk-Forward 검증'),
    (C_SKY,  '라벨 설계',
     '유지 37% → 노이즈\n이진 분류 전환\nprice[T+h] vs price[T]'),
    (C_BLUE, '피처 설계',
     '선행지표 Lag 피처\n변화율·모멘텀\nMA 크로스 신호'),
    (C_NAVY, '최종 결과',
     '10일 Accuracy\n44% → 78.9%\nAUC 0.852'),
]
fw = Inches(2.35)
fh = Inches(1.7)
fy = Inches(0.88)
fx = Inches(0.3)
for i, (col, ftitle, fbody) in enumerate(flow_items):
    box(s3, fx, fy, fw, fh, fill=col, rnd=True)
    txt(s3, ftitle, fx, fy + Inches(0.12), fw, Inches(0.38),
        size=Pt(10.5), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s3, fbody,
        fx + Inches(0.12), fy + Inches(0.55),
        fw - Inches(0.24), Inches(1.05),
        size=Pt(9), color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        txt(s3, '▶',
            fx + fw + Inches(0.02), fy + Inches(0.6),
            Inches(0.3), Inches(0.5),
            size=Pt(16), color=C_LIGHT, bold=True, align=PP_ALIGN.CENTER)
    fx += fw + Inches(0.32)

# ── 좌: 예측 타임라인 ────────────────────────────────────
box(s3, Inches(0.25), Inches(2.75), Inches(7.9), Inches(4.05),
    fill=C_WHITE, line=C_LIGHT, lw=Pt(1.2), rnd=True)
txt(s3, '예측 방향 시계열  —  Ensemble 모델 (30일 기준)',
    Inches(0.4), Inches(2.84), Inches(7.4), Inches(0.32),
    size=Pt(10), color=C_BLUE, bold=True)
img(s3, '08_pred_timeline.png',
    Inches(0.3), Inches(3.18), Inches(7.8), Inches(3.3))

# ── 우: Before/After 수치 카드 ───────────────────────────
bfx = Inches(8.35)
ba_data = [
    ('10일', '44.1%', '78.9%', '0.852', C_BLUE),
    ('30일', '36.8%', '52.7%', '0.728', C_SKY),
    ('60일', '39.6%', '53.4%', '0.741', C_LIGHT),
]
txt(s3, 'Accuracy  Before → After',
    bfx, Inches(2.84), Inches(4.7), Inches(0.32),
    size=Pt(10), color=C_BLUE, bold=True)

for i, (period, before, after, auc, col) in enumerate(ba_data):
    cy = Inches(3.22) + i * Inches(1.28)
    cw = Inches(4.72)
    box(s3, bfx, cy, cw, Inches(1.15),
        fill=C_WHITE, line=col, lw=Pt(1.8), rnd=True)
    box(s3, bfx, cy, Inches(0.08), Inches(1.15), fill=col)
    # 기간
    txt(s3, period,
        bfx + Inches(0.18), cy + Inches(0.1),
        Inches(0.7), Inches(0.35),
        size=Pt(13), color=col, bold=True)
    # Before
    txt(s3, before,
        bfx + Inches(0.95), cy + Inches(0.08),
        Inches(1.1), Inches(0.46),
        size=Pt(18), color=C_GRAY, bold=True, align=PP_ALIGN.CENTER)
    txt(s3, 'BEFORE',
        bfx + Inches(0.95), cy + Inches(0.56),
        Inches(1.1), Inches(0.24),
        size=Pt(7.5), color=C_GRAY, align=PP_ALIGN.CENTER)
    # 화살표
    txt(s3, '→',
        bfx + Inches(2.05), cy + Inches(0.12),
        Inches(0.5), Inches(0.46),
        size=Pt(20), color=col, bold=True, align=PP_ALIGN.CENTER)
    # After
    txt(s3, after,
        bfx + Inches(2.55), cy + Inches(0.08),
        Inches(1.1), Inches(0.46),
        size=Pt(18), color=col, bold=True, align=PP_ALIGN.CENTER)
    txt(s3, 'AFTER',
        bfx + Inches(2.55), cy + Inches(0.56),
        Inches(1.1), Inches(0.24),
        size=Pt(7.5), color=col, align=PP_ALIGN.CENTER)
    # AUC
    txt(s3, f'AUC {auc}',
        bfx + Inches(3.72), cy + Inches(0.35),
        Inches(0.85), Inches(0.28),
        size=Pt(8.5), color=col, italic=True, align=PP_ALIGN.CENTER)

# 하단 결론 바
box(s3, Inches(0), Inches(7.08), W, Inches(0.04), fill=C_SKY)
txt(s3,
    '데이터 분석 인사이트(비정상·선행지표·유지 노이즈)를 그대로 모델 설계에 반영  →  10일 예측 정확도 +34.8%p 향상',
    Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3),
    size=Pt(9.5), color=C_BLUE, bold=True, align=PP_ALIGN.CENTER)


prs.save('outputs/analysis_ppt_v2.pptx')
print(f'저장 완료: outputs/analysis_ppt_v2.pptx  ({len(prs.slides)}슬라이드)')

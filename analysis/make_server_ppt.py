# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 ──────────────────────────────────────────────────────
C_BG       = RGBColor(0xF4, 0xF6, 0xF9)   # 연회색 배경
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY     = RGBColor(0x6C, 0x75, 0x7D)
C_LGRAY    = RGBColor(0xDE, 0xE2, 0xE6)

# 섹션 컬러
C_NODE     = RGBColor(0x68, 0xB7, 0x2A)   # Node.js 초록
C_FLASK    = RGBColor(0x00, 0x7B, 0xFF)   # Flask 파랑
C_DB       = RGBColor(0xF7, 0x7F, 0x00)   # DB 주황
C_CRAWL    = RGBColor(0x8A, 0x2B, 0xE2)   # 크롤러 보라
C_MODEL    = RGBColor(0xE8, 0x3E, 0x3E)   # 모델 빨강
C_ARROW    = RGBColor(0x49, 0x50, 0x57)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 배경
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = C_BG


# ── 헬퍼 ──────────────────────────────────────────────────────
def box(x, y, w, h, fill=C_WHITE, line=None, lw=Pt(1), rnd=False):
    s = slide.shapes.add_shape(5 if rnd else 1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    if rnd and hasattr(s, 'adjustments'):
        s.adjustments[0] = 0.08
    return s

def txt(text, x, y, w, h, size=Pt(9), color=C_DARK,
        bold=False, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = size
    r.font.color.rgb = color
    r.font.bold   = bold
    r.font.italic = italic
    r.font.name   = 'Malgun Gothic'
    return tb

def arrow_h(x1, y, x2, color=C_ARROW, w=Pt(2.5)):
    c = slide.shapes.add_connector(1, x1, y, x2, y)
    c.line.color.rgb = color
    c.line.width = w
    return c

def arrow_v(x, y1, y2, color=C_ARROW, w=Pt(1.5)):
    c = slide.shapes.add_connector(1, x, y1, x, y2)
    c.line.color.rgb = color
    c.line.width = w
    return c

def tag(text, x, y, w, h, fill, text_color=C_WHITE):
    box(x, y, w, h, fill=fill, rnd=True)
    txt(text, x, y, w, h, size=Pt(8), color=text_color,
        bold=True, align=PP_ALIGN.CENTER)

def mini_box(text, x, y, w, h, line_color, fill=C_WHITE):
    box(x, y, w, h, fill=fill, line=line_color, lw=Pt(1.2), rnd=True)
    txt(text, x + Inches(0.08), y + Inches(0.04),
        w - Inches(0.16), h - Inches(0.08),
        size=Pt(8.5), color=C_DARK, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════
# 제목
# ════════════════════════════════════════════════════════════════
box(Inches(0), Inches(0), W, Inches(0.72), fill=C_DARK)
txt('서버 아키텍처', Inches(0.4), Inches(0.1), Inches(4), Inches(0.52),
    size=Pt(22), color=C_WHITE, bold=True, align=PP_ALIGN.LEFT)
txt('Browser  →  Node.js (Web Server)  →  Flask (Backend)  →  MySQL DB',
    Inches(4.5), Inches(0.18), Inches(8.5), Inches(0.38),
    size=Pt(10), color=RGBColor(0xAA, 0xBB, 0xCC),
    align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════
# 섹션 1 : 브라우저  (좌)
# ════════════════════════════════════════════════════════════════
S1_X = Inches(0.25)
S1_Y = Inches(0.85)
S1_W = Inches(2.2)
S1_H = Inches(6.3)

box(S1_X, S1_Y, S1_W, S1_H,
    fill=RGBColor(0xEA, 0xF6, 0xFF),
    line=RGBColor(0x90, 0xC8, 0xF0), lw=Pt(1.5))

# 섹션 제목 바
box(S1_X, S1_Y, S1_W, Inches(0.5),
    fill=RGBColor(0x34, 0x8A, 0xD4), line=None)
txt('Browser / Client', S1_X, S1_Y + Inches(0.08), S1_W, Inches(0.38),
    size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# 페이지 카드들
pages = [
    ('index.html',  '대시보드 메인'),
    ('chart.html',  '예측 차트'),
    ('tables.html', '단가 테이블'),
    ('craw.html',   '크롤링 현황'),
]
for i, (fname, fdesc) in enumerate(pages):
    py = S1_Y + Inches(0.7) + i * Inches(1.12)
    box(S1_X + Inches(0.18), py, S1_W - Inches(0.36), Inches(0.95),
        fill=C_WHITE, line=RGBColor(0x90, 0xC8, 0xF0), lw=Pt(1.2), rnd=True)
    # 브라우저 탭 흉내
    box(S1_X + Inches(0.18), py, S1_W - Inches(0.36), Inches(0.22),
        fill=RGBColor(0xD0, 0xE8, 0xF8), line=None)
    txt('⬛ ⬛ ⬛', S1_X + Inches(0.25), py + Inches(0.03),
        Inches(0.6), Inches(0.18), size=Pt(6), color=C_GRAY)
    txt(fname, S1_X + Inches(0.26), py + Inches(0.28),
        S1_W - Inches(0.5), Inches(0.3),
        size=Pt(9.5), color=RGBColor(0x34, 0x8A, 0xD4), bold=True)
    txt(fdesc, S1_X + Inches(0.26), py + Inches(0.58),
        S1_W - Inches(0.5), Inches(0.28),
        size=Pt(8), color=C_GRAY)

# 하단 포트 표시
txt('Port: 8080', S1_X, S1_Y + S1_H - Inches(0.35),
    S1_W, Inches(0.3),
    size=Pt(8.5), color=RGBColor(0x34, 0x8A, 0xD4),
    bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 화살표 1: 브라우저 → Node
# ════════════════════════════════════════════════════════════════
ARR1_Y = Inches(3.8)
arrow_h(S1_X + S1_W, ARR1_Y, S1_X + S1_W + Inches(0.38),
        color=RGBColor(0x34, 0x8A, 0xD4), w=Pt(3))
txt('HTTP', S1_X + S1_W + Inches(0.02), ARR1_Y - Inches(0.28),
    Inches(0.5), Inches(0.24), size=Pt(7.5),
    color=RGBColor(0x34, 0x8A, 0xD4), bold=True)


# ════════════════════════════════════════════════════════════════
# 섹션 2 : Node.js 웹 서버
# ════════════════════════════════════════════════════════════════
S2_X = S1_X + S1_W + Inches(0.38)
S2_Y = Inches(0.85)
S2_W = Inches(3.0)
S2_H = Inches(6.3)

box(S2_X, S2_Y, S2_W, S2_H,
    fill=RGBColor(0xF0, 0xFB, 0xEA),
    line=RGBColor(0x90, 0xD0, 0x80), lw=Pt(1.5))

box(S2_X, S2_Y, S2_W, Inches(0.5),
    fill=C_NODE, line=None)
txt('Node.js  (Express)', S2_X, S2_Y + Inches(0.08), S2_W, Inches(0.38),
    size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

txt('server.js', S2_X + Inches(0.15), S2_Y + Inches(0.58),
    Inches(1.5), Inches(0.28),
    size=Pt(8.5), color=C_NODE, bold=True)

# 역할 박스
roles = [
    '정적 파일 서빙\n(HTML / CSS / JS)',
    'MySQL 직접 조회\n(단가·피처 데이터)',
    'Flask API 중계\n(예측·크롤링 요청)',
    'CORS 처리\n환경변수 설정 주입',
]
for i, r in enumerate(roles):
    ry = S2_Y + Inches(0.98) + i * Inches(1.12)
    box(S2_X + Inches(0.15), ry,
        S2_W - Inches(0.3), Inches(0.95),
        fill=C_WHITE, line=C_NODE, lw=Pt(1.2), rnd=True)
    # 좌측 바
    box(S2_X + Inches(0.15), ry, Inches(0.07), Inches(0.95),
        fill=C_NODE, line=None)
    txt(r, S2_X + Inches(0.3), ry + Inches(0.1),
        S2_W - Inches(0.5), Inches(0.78),
        size=Pt(9), color=C_DARK)

txt('Port: 8080', S2_X, S2_Y + S2_H - Inches(0.35),
    S2_W, Inches(0.3),
    size=Pt(8.5), color=C_NODE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 화살표 2: Node → Flask
# ════════════════════════════════════════════════════════════════
ARR2_Y = Inches(3.8)
arrow_h(S2_X + S2_W, ARR2_Y, S2_X + S2_W + Inches(0.38),
        color=C_FLASK, w=Pt(3))
txt('REST', S2_X + S2_W + Inches(0.02), ARR2_Y - Inches(0.28),
    Inches(0.5), Inches(0.24), size=Pt(7.5),
    color=C_FLASK, bold=True)


# ════════════════════════════════════════════════════════════════
# 섹션 3 : Flask 백엔드
# ════════════════════════════════════════════════════════════════
S3_X = S2_X + S2_W + Inches(0.38)
S3_Y = Inches(0.85)
S3_W = Inches(3.8)
S3_H = Inches(6.3)

box(S3_X, S3_Y, S3_W, S3_H,
    fill=RGBColor(0xEB, 0xF2, 0xFF),
    line=RGBColor(0x80, 0xA8, 0xF0), lw=Pt(1.5))

box(S3_X, S3_Y, S3_W, Inches(0.5),
    fill=C_FLASK, line=None)
txt('Flask  (Python 3.11)', S3_X, S3_Y + Inches(0.08), S3_W, Inches(0.38),
    size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

txt('app_test.py  |  predict.py  |  crawlers/',
    S3_X + Inches(0.15), S3_Y + Inches(0.58),
    S3_W - Inches(0.3), Inches(0.28),
    size=Pt(8), color=C_FLASK, bold=True)

# API 엔드포인트 그룹
api_groups = [
    (C_CRAWL, '크롤링',
     '/run-crawling\n/save-data'),
    (C_MODEL, '예측',
     '/transfer-predict\n/run-all-versions'),
    (RGBColor(0x00, 0x96, 0x88), '조회',
     '/get-chart-data\n/get-dashboard-data\n/get-recent-data'),
    (C_GRAY,  '기타',
     '/interpolate-data\n/download-data-file'),
]
for i, (col, gname, apis) in enumerate(api_groups):
    gy = S3_Y + Inches(1.0) + i * Inches(1.2)
    box(S3_X + Inches(0.15), gy,
        S3_W - Inches(0.3), Inches(1.05),
        fill=C_WHITE, line=col, lw=Pt(1.2), rnd=True)
    box(S3_X + Inches(0.15), gy, Inches(0.07), Inches(1.05),
        fill=col, line=None)
    # 그룹 이름 배지
    box(S3_X + Inches(0.32), gy + Inches(0.1),
        Inches(0.7), Inches(0.26),
        fill=col, line=None, rnd=True)
    txt(gname, S3_X + Inches(0.32), gy + Inches(0.1),
        Inches(0.7), Inches(0.26),
        size=Pt(7.5), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(apis, S3_X + Inches(0.32), gy + Inches(0.42),
        S3_W - Inches(0.55), Inches(0.55),
        size=Pt(8.5), color=C_DARK)

txt('Port: 5050', S3_X, S3_Y + S3_H - Inches(0.35),
    S3_W, Inches(0.3),
    size=Pt(8.5), color=C_FLASK, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 화살표 3: Flask ↔ DB
# ════════════════════════════════════════════════════════════════
ARR3_Y = Inches(3.8)
arrow_h(S3_X + S3_W, ARR3_Y, S3_X + S3_W + Inches(0.38),
        color=C_DB, w=Pt(3))
txt('SQL', S3_X + S3_W + Inches(0.03), ARR3_Y - Inches(0.28),
    Inches(0.45), Inches(0.24), size=Pt(7.5),
    color=C_DB, bold=True)


# ════════════════════════════════════════════════════════════════
# 섹션 4 : MySQL DB
# ════════════════════════════════════════════════════════════════
S4_X = S3_X + S3_W + Inches(0.38)
S4_Y = Inches(0.85)
S4_W = Inches(2.55)
S4_H = Inches(6.3)

box(S4_X, S4_Y, S4_W, S4_H,
    fill=RGBColor(0xFF, 0xF6, 0xE8),
    line=RGBColor(0xF0, 0xA0, 0x30), lw=Pt(1.5))

box(S4_X, S4_Y, S4_W, Inches(0.5),
    fill=C_DB, line=None)
txt('MySQL  (steel_db)', S4_X, S4_Y + Inches(0.08), S4_W, Inches(0.38),
    size=Pt(11), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

tables = [
    ('steel_prices',    '일별 단가·물동량\n2,663행'),
    ('steel_features',  'EAV 피처 시계열\n386,135행 / 155종'),
    ('predictions',     '예측 결과 저장\nbase_date · horizon'),
    ('crawl_logs',      '크롤링 실행 이력'),
]
for i, (tname, tdesc) in enumerate(tables):
    ty = S4_Y + Inches(0.68) + i * Inches(1.32)
    box(S4_X + Inches(0.15), ty,
        S4_W - Inches(0.3), Inches(1.15),
        fill=C_WHITE, line=C_DB, lw=Pt(1.2), rnd=True)

    # 실린더 아이콘 흉내 (상단 타원)
    box(S4_X + Inches(0.22), ty + Inches(0.06),
        Inches(0.32), Inches(0.22),
        fill=RGBColor(0xFF, 0xDD, 0xAA), line=None, rnd=True)
    box(S4_X + Inches(0.22), ty + Inches(0.17),
        Inches(0.32), Inches(0.14),
        fill=RGBColor(0xFF, 0xDD, 0xAA), line=None)

    txt(tname, S4_X + Inches(0.65), ty + Inches(0.1),
        S4_W - Inches(0.85), Inches(0.3),
        size=Pt(9), color=C_DB, bold=True)
    txt(tdesc, S4_X + Inches(0.65), ty + Inches(0.45),
        S4_W - Inches(0.85), Inches(0.6),
        size=Pt(8), color=C_GRAY)


# ════════════════════════════════════════════════════════════════
# 하단 크롤러 소스 바
# ════════════════════════════════════════════════════════════════
box(Inches(0.25), Inches(7.08), W - Inches(0.5), Inches(0.35),
    fill=RGBColor(0xF0, 0xE8, 0xFF), line=C_CRAWL, lw=Pt(1))

txt('크롤링 소스 :', Inches(0.35), Inches(7.1), Inches(1.2), Inches(0.28),
    size=Pt(8.5), color=C_CRAWL, bold=True)

sources = ['steelprice.co.kr  (단가)', 'steeldaily.co.kr  (시황)',
           'steelin.co.kr  (물동량)', 'ECOS API  (한국은행)', 'KOSIS API  (통계청)']
sx = Inches(1.55)
for s in sources:
    box(sx, Inches(7.11), Inches(2.15), Inches(0.26),
        fill=RGBColor(0xDE, 0xCC, 0xFF), line=C_CRAWL, lw=Pt(0.8), rnd=True)
    txt(s, sx + Inches(0.08), Inches(7.11), Inches(2.0), Inches(0.26),
        size=Pt(7.5), color=C_CRAWL)
    sx += Inches(2.3)

prs.save('outputs/server_slide.pptx')
print('저장 완료: outputs/server_slide.pptx')
